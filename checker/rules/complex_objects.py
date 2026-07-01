"""Manual-review checks for complex objects."""

from __future__ import annotations

from typing import Any

from pptx.enum.shapes import MSO_SHAPE_TYPE

from config import Settings
from models import Finding, RiskLevel
from utils.pptx_xml import iter_shapes


def _table_dimensions(shape: Any) -> tuple[int, int]:
    """Return row and column count for a table shape."""
    try:
        return len(shape.table.rows), len(shape.table.columns)
    except (AttributeError, ValueError):
        return 0, 0


def check(prs: Any, settings: Settings) -> list[Finding]:
    """Flag charts, media, and large tables for human review."""
    del settings
    findings: list[Finding] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in iter_shapes(slide.shapes):
            shape_id = str(shape.shape_id)
            if getattr(shape, "has_chart", False):
                findings.append(
                    Finding(
                        rule_id="chart_requires_manual_alt_text",
                        slide_number=slide_index,
                        element_id=shape_id,
                        element_type="chart",
                        wcag_criterion="1.1.1 Non-text Content",
                        section_508_ref="E205.4",
                        risk_level=RiskLevel.HIGH,
                        issue_description=(
                            "Chart or data graphic needs a meaningful text "
                            "alternative. This requires human interpretation of "
                            "the data and message."
                        ),
                        suggested_fix=None,
                        metadata={},
                    )
                )
            if getattr(shape, "has_table", False):
                rows, columns = _table_dimensions(shape)
                if rows > 8 or columns > 6:
                    findings.append(
                        Finding(
                            rule_id="complex_table_review",
                            slide_number=slide_index,
                            element_id=shape_id,
                            element_type="table",
                            wcag_criterion="1.3.1 Info and Relationships",
                            section_508_ref="E205.4",
                            risk_level=RiskLevel.HIGH,
                            issue_description=(
                                f"Large table detected ({rows} rows x {columns} columns). "
                                "Confirm headers, reading order, and whether a summary is needed."
                            ),
                            suggested_fix=None,
                            metadata={"rows": rows, "columns": columns},
                        )
                    )
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                findings.append(
                    Finding(
                        rule_id="media_requires_captions",
                        slide_number=slide_index,
                        element_id=shape_id,
                        element_type="media",
                        wcag_criterion="1.2.2 Captions (Prerecorded)",
                        section_508_ref="E205.4",
                        risk_level=RiskLevel.HIGH,
                        issue_description=(
                            "Embedded audio or video needs captions, transcript, "
                            "or another accessible alternative."
                        ),
                        suggested_fix=None,
                        metadata={},
                    )
                )
    return findings
