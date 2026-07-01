"""Image alternative text accessibility checks."""

from __future__ import annotations

import re
from typing import Any

from pptx.enum.shapes import MSO_SHAPE_TYPE

from config import Settings
from models import Finding, RiskLevel
from utils.pptx_xml import get_alt_text, iter_shapes, slide_context_text, title_text

FILENAME_PATTERN = re.compile(r"^(image|img|picture|photo|diagram)[-_ ]?\d*\.(png|jpe?g|gif|bmp|svg)$", re.I)
AUTOGEN_PATTERN = re.compile(r"^(image|img|picture|photo|graphic)[-_ ]?\d+$", re.I)
EXTENSION_PATTERN = re.compile(r"^[\w .-]{1,120}\.(png|jpe?g|gif|bmp|svg|webp|tiff?)$", re.I)


def looks_like_weak_alt_text(text: str) -> bool:
    """Return True when alt text appears generated or non-descriptive."""
    stripped = text.strip()
    if not stripped:
        return True
    if FILENAME_PATTERN.match(stripped) or AUTOGEN_PATTERN.match(stripped) or EXTENSION_PATTERN.match(stripped):
        return True
    return stripped.lower() in {"image", "picture", "photo", "graphic", "diagram"}


def _suggest_alt_text(slide: Any, shape_id: str, slide_number: int) -> str:
    """Generate a human-reviewable alt-text suggestion from slide context."""
    title = title_text(slide)
    context = slide_context_text(slide, exclude_shape_id=shape_id)
    if title and context:
        return f"Image related to '{title}': {context}"
    if title:
        return f"Image related to '{title}'"
    if context:
        return f"Image related to: {context}"
    return f"Describe the purpose of this image on slide {slide_number}."


def check(prs: Any, settings: Settings) -> list[Finding]:
    """Detect missing or weak alt text on picture shapes."""
    del settings
    findings: list[Finding] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in iter_shapes(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            shape_id = str(shape.shape_id)
            alt_text = get_alt_text(shape)
            if not alt_text:
                findings.append(
                    Finding(
                        rule_id="missing_image_alt_text",
                        slide_number=slide_index,
                        element_id=shape_id,
                        element_type="image",
                        wcag_criterion="1.1.1 Non-text Content",
                        section_508_ref="E205.4",
                        risk_level=RiskLevel.MEDIUM,
                        issue_description=(
                            "Image has no alt text. A human should approve or "
                            "edit the suggested description before it is applied."
                        ),
                        suggested_fix=_suggest_alt_text(slide, shape_id, slide_index),
                        metadata={"before_alt_text": alt_text},
                    )
                )
            elif looks_like_weak_alt_text(alt_text):
                findings.append(
                    Finding(
                        rule_id="weak_image_alt_text",
                        slide_number=slide_index,
                        element_id=shape_id,
                        element_type="image",
                        wcag_criterion="1.1.1 Non-text Content",
                        section_508_ref="E205.4",
                        risk_level=RiskLevel.MEDIUM,
                        issue_description=(
                            f"Image alt text '{alt_text}' appears generic, "
                            "filename-like, or automatically generated."
                        ),
                        suggested_fix=_suggest_alt_text(slide, shape_id, slide_index),
                        metadata={"before_alt_text": alt_text},
                    )
                )
    return findings
