"""End-to-end tests for the local web app via FastAPI's TestClient."""

from __future__ import annotations

import base64
import io
from pathlib import Path

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

from web.jobs import JobStore
from web.server import create_app

ONE_PIXEL_PNG = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p94AAAAASUVORK5CYII="


def _build_deck(tmp_path: Path) -> Path:
    """Create a deck with low-risk (link/title/language) and medium-risk (alt-text) issues."""
    image_path = tmp_path / "image001.png"
    image_path.write_bytes(base64.b64decode(ONE_PIXEL_PNG))
    pptx_path = tmp_path / "deck.pptx"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "click here"
    run.hyperlink.address = "https://example.com/resource"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(3), width=Inches(1))
    prs.save(pptx_path)
    return pptx_path


def _client(tmp_path: Path) -> TestClient:
    app = create_app(store=JobStore.create(tmp_path / "jobs"))
    return TestClient(app)


def _analyze(client: TestClient, pptx_path: Path) -> dict:
    with pptx_path.open("rb") as handle:
        res = client.post(
            "/api/analyze",
            files={"file": ("deck.pptx", handle, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            data={"api_key": "", "provider": "anthropic", "reviewer": "tester"},
        )
    assert res.status_code == 200, res.text
    return res.json()


def test_index_served(tmp_path: Path) -> None:
    """The root serves the single-page app."""
    res = _client(tmp_path).get("/")
    assert res.status_code == 200
    assert "AccessiSlides" in res.text


def test_analyze_returns_job_and_ledger(tmp_path: Path) -> None:
    """Analyze auto-fixes low-risk items and stages medium-risk ones."""
    client = _client(tmp_path)
    data = _analyze(client, _build_deck(tmp_path))
    assert data["job_id"]
    assert data["llm_active"] is False
    assert data["counts"]["auto"] >= 1
    assert data["counts"]["pending"] >= 1
    statuses = {item["status"] for item in data["items"]}
    assert "auto_applied" in statuses
    assert "pending_approval" in statuses


def test_approve_and_reject_transitions(tmp_path: Path) -> None:
    """A pending suggestion can be approved; another can be rejected."""
    client = _client(tmp_path)
    data = _analyze(client, _build_deck(tmp_path))
    job_id = data["job_id"]
    pending = [it["item_id"] for it in data["items"] if it["status"] == "pending_approval"]
    assert pending

    approved = client.post(f"/api/jobs/{job_id}/items/{pending[0]}/approve").json()
    target = next(it for it in approved["items"] if it["item_id"] == pending[0])
    assert target["status"] == "approved"
    assert target["approved_by"] == "human:tester"

    rejected = client.post(f"/api/jobs/{job_id}/items/{pending[0]}/reject")
    # already approved → cannot reject
    assert rejected.status_code == 409


def test_rollback_auto_applied(tmp_path: Path) -> None:
    """An auto-applied low-risk fix can be rolled back."""
    client = _client(tmp_path)
    data = _analyze(client, _build_deck(tmp_path))
    job_id = data["job_id"]
    auto_item = next(it["item_id"] for it in data["items"] if it["status"] == "auto_applied")
    rolled = client.post(f"/api/jobs/{job_id}/items/{auto_item}/rollback").json()
    target = next(it for it in rolled["items"] if it["item_id"] == auto_item)
    assert target["status"] == "rolled_back"


def test_download_returns_valid_pptx(tmp_path: Path) -> None:
    """The download endpoint returns an openable presentation."""
    client = _client(tmp_path)
    data = _analyze(client, _build_deck(tmp_path))
    res = client.get(f"/api/jobs/{data['job_id']}/download")
    assert res.status_code == 200
    prs = Presentation(io.BytesIO(res.content))
    assert len(prs.slides) == 1


def test_export_xlsx_and_pdf(tmp_path: Path) -> None:
    """Audit reports export in both formats."""
    client = _client(tmp_path)
    job_id = _analyze(client, _build_deck(tmp_path))["job_id"]
    xlsx = client.get(f"/api/jobs/{job_id}/export", params={"format": "xlsx"})
    pdf = client.get(f"/api/jobs/{job_id}/export", params={"format": "pdf"})
    assert xlsx.status_code == 200 and xlsx.content[:2] == b"PK"
    assert pdf.status_code == 200 and pdf.content[:4] == b"%PDF"


def test_unknown_job_is_404(tmp_path: Path) -> None:
    """Operations on an unknown job return 404."""
    res = _client(tmp_path).get("/api/jobs/nope/ledger")
    assert res.status_code == 404


def _analyze_with(client: TestClient, pptx_path: Path, data: dict):
    """POST /api/analyze with explicit form fields, returning the raw response."""
    with pptx_path.open("rb") as handle:
        return client.post(
            "/api/analyze",
            files={"file": ("deck.pptx", handle, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            data=data,
        )


def test_llm_enabled_without_key_is_400(tmp_path: Path) -> None:
    """Turning AI on without an API key yields a clear error (popup on the client)."""
    client = _client(tmp_path)
    res = _analyze_with(client, _build_deck(tmp_path), {"use_llm": "true", "api_key": "", "provider": "anthropic"})
    assert res.status_code == 400
    assert "API key" in res.json()["detail"]


def test_llm_enabled_unknown_provider_is_400(tmp_path: Path) -> None:
    """An unsupported provider is rejected before any analysis runs."""
    client = _client(tmp_path)
    res = _analyze_with(client, _build_deck(tmp_path), {"use_llm": "true", "api_key": "x", "provider": "acme"})
    assert res.status_code == 400
    assert "provider" in res.json()["detail"].lower()


def test_llm_provider_error_is_502(tmp_path: Path, monkeypatch) -> None:
    """A provider/auth/model failure during preflight surfaces as a 502 error."""

    class BoomProvider:
        name = "anthropic:bad-model"

        def generate_text(self, *, system: str, prompt: str) -> str:
            raise RuntimeError("401 invalid api key")

        def describe_image(self, **_: object) -> str:  # pragma: no cover - unused
            return ""

    monkeypatch.setattr("llm.factory.get_provider", lambda settings: BoomProvider())
    client = _client(tmp_path)
    res = _analyze_with(
        client, _build_deck(tmp_path),
        {"use_llm": "true", "api_key": "x", "provider": "anthropic", "model": "bad-model"},
    )
    assert res.status_code == 502
    assert "invalid api key" in res.json()["detail"]


def test_llm_disabled_ignores_key_and_model(tmp_path: Path) -> None:
    """With AI off, the deterministic pipeline runs even if a key/model are sent."""
    client = _client(tmp_path)
    res = _analyze_with(
        client, _build_deck(tmp_path),
        {"use_llm": "false", "api_key": "x", "provider": "anthropic", "model": "whatever"},
    )
    assert res.status_code == 200
    assert res.json()["llm_active"] is False


def test_env_status_reports_saved_keys_without_leaking(tmp_path: Path, monkeypatch) -> None:
    """The env endpoint flags a saved key by boolean and never returns the value."""
    monkeypatch.setenv("PPTXA_ANTHROPIC_API_KEY", "env-secret-value")
    res = _client(tmp_path).get("/api/env")
    assert res.status_code == 200
    body = res.json()
    assert body["saved_keys"]["anthropic"] is True
    assert "openai" in body["saved_keys"]
    assert "provider" in body
    assert "env-secret-value" not in res.text


def test_saved_key_uses_env_key(tmp_path: Path, monkeypatch) -> None:
    """With use_saved_key on, the provider is built from the .env key, not the form."""
    monkeypatch.setenv("PPTXA_ANTHROPIC_API_KEY", "env-secret-value")

    class OkProvider:
        name = "anthropic:claude"

        def generate_text(self, *, system: str, prompt: str) -> str:
            return "OK"

        def describe_image(self, **_: object) -> str:  # pragma: no cover - unused
            return ""

    captured: dict[str, object] = {}

    def fake_get_provider(settings):
        captured["key"] = settings.anthropic_api_key
        return OkProvider()

    monkeypatch.setattr("llm.factory.get_provider", fake_get_provider)
    client = _client(tmp_path)
    res = _analyze_with(
        client, _build_deck(tmp_path),
        {"use_llm": "true", "use_saved_key": "true", "provider": "anthropic", "api_key": ""},
    )
    assert res.status_code == 200, res.text
    assert res.json()["llm_active"] is True
    assert captured["key"] == "env-secret-value"


def test_saved_key_missing_is_400(tmp_path: Path, monkeypatch) -> None:
    """Asking to use a saved key that is not set yields a clear API-key error."""
    monkeypatch.setenv("PPTXA_ANTHROPIC_API_KEY", "")
    client = _client(tmp_path)
    res = _analyze_with(
        client, _build_deck(tmp_path),
        {"use_llm": "true", "use_saved_key": "true", "provider": "anthropic", "api_key": ""},
    )
    assert res.status_code == 400
    assert "API key" in res.json()["detail"]


def test_bad_pptx_returns_document_error(tmp_path: Path) -> None:
    """An unreadable .pptx yields a 400 the client shows as a document-error popup."""
    bad = tmp_path / "broken.pptx"
    bad.write_bytes(b"not a real pptx")
    client = _client(tmp_path)
    res = _analyze_with(client, bad, {"use_llm": "false"})
    assert res.status_code == 400
    assert "Could not read presentation" in res.json()["detail"]


def _isolated_env(tmp_path: Path, monkeypatch) -> Path:
    """Point the keystore at a throwaway .env and clear any exported keys."""
    import keystore

    env_path = tmp_path / ".env"
    monkeypatch.setattr(keystore, "ENV_PATH", env_path)
    monkeypatch.delenv("PPTXA_ANTHROPIC_API_KEY", raising=False)
    monkeypatch.delenv("PPTXA_OPENAI_API_KEY", raising=False)
    return env_path


def test_save_key_endpoint_persists_key_and_model(tmp_path: Path, monkeypatch) -> None:
    """Saving stores the key plus the model chosen with it, and never echoes the key."""
    env_path = _isolated_env(tmp_path, monkeypatch)
    client = _client(tmp_path)

    res = client.post(
        "/api/key/save",
        json={"provider": "anthropic", "api_key": "sk-ant-x", "model": "claude-sonnet-5"},
    )
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["saved_keys"]["anthropic"] is True
    assert body["key_sources"]["anthropic"] == "env_file"
    assert body["saved_models"]["anthropic"] == "claude-sonnet-5"
    assert body["erasable"] == ["anthropic"]
    assert "sk-ant-x" not in res.text  # the value itself never comes back
    text = env_path.read_text(encoding="utf-8")
    assert "PPTXA_ANTHROPIC_API_KEY=sk-ant-x" in text
    assert "PPTXA_ANTHROPIC_MODEL=claude-sonnet-5" in text


def test_save_key_endpoint_keeps_both_providers(tmp_path: Path, monkeypatch) -> None:
    """Both keys can be saved at once; the second save does not evict the first."""
    _isolated_env(tmp_path, monkeypatch)
    client = _client(tmp_path)
    client.post("/api/key/save", json={"provider": "openai", "api_key": "sk-o", "model": "gpt-4o"})
    res = client.post(
        "/api/key/save", json={"provider": "anthropic", "api_key": "sk-a", "model": "claude-sonnet-5"}
    )

    body = res.json()
    assert body["saved_keys"] == {"anthropic": True, "openai": True}
    assert body["saved_models"] == {"anthropic": "claude-sonnet-5", "openai": "gpt-4o"}
    assert sorted(body["erasable"]) == ["anthropic", "openai"]


def test_env_provider_is_not_guessed_when_both_keys_saved(tmp_path: Path, monkeypatch) -> None:
    """With one key the UI preselects it; with two it falls back to the configured default."""
    _isolated_env(tmp_path, monkeypatch)
    client = _client(tmp_path)

    client.post("/api/key/save", json={"provider": "anthropic", "api_key": "sk-a"})
    assert client.get("/api/env").json()["provider"] == "anthropic"

    client.post("/api/key/save", json={"provider": "openai", "api_key": "sk-o"})
    assert client.get("/api/env").json()["provider"] == "openai"  # config.py default


def test_save_key_endpoint_rejects_bad_provider(tmp_path: Path, monkeypatch) -> None:
    """An unsupported provider is a 400, and nothing is written."""
    env_path = _isolated_env(tmp_path, monkeypatch)
    res = _client(tmp_path).post("/api/key/save", json={"provider": "acme", "api_key": "sk-x"})
    assert res.status_code == 400
    assert not env_path.exists()


def test_forget_key_endpoint_erases_one_provider(tmp_path: Path, monkeypatch) -> None:
    """Erasing a named provider leaves the other provider's key saved."""
    env_path = _isolated_env(tmp_path, monkeypatch)
    client = _client(tmp_path)
    client.post("/api/key/save", json={"provider": "openai", "api_key": "sk-openai-x"})
    client.post("/api/key/save", json={"provider": "anthropic", "api_key": "sk-ant-x"})

    res = client.post("/api/key/forget", json={"providers": ["anthropic"]})
    assert res.status_code == 200, res.text
    body = res.json()
    assert body["removed"] == ["anthropic"]
    assert body["saved_keys"] == {"anthropic": False, "openai": True}
    text = env_path.read_text(encoding="utf-8")
    assert "sk-ant-x" not in text
    assert "sk-openai-x" in text


def test_forget_key_endpoint_erases_all_by_default(tmp_path: Path, monkeypatch) -> None:
    """A bodyless forget clears every saved key."""
    env_path = _isolated_env(tmp_path, monkeypatch)
    client = _client(tmp_path)
    client.post("/api/key/save", json={"provider": "openai", "api_key": "sk-openai-x"})
    client.post("/api/key/save", json={"provider": "anthropic", "api_key": "sk-ant-x"})

    body = client.post("/api/key/forget").json()
    assert sorted(body["removed"]) == ["anthropic", "openai"]
    assert body["can_erase"] is False
    assert "sk-openai-x" not in env_path.read_text(encoding="utf-8")
