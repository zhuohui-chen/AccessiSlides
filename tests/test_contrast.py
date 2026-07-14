"""Tests for the WCAG 1.4.3 contrast rule and its color math.

The rule is fully deterministic: theme and inherited colors are resolved to
concrete RGB via the slide master's palette, so no LLM is involved.
"""

from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Inches, Pt

from checker.engine import run_checks_on_presentation
from checker.rules import contrast
from checker.triage import classify
from config import Settings
from models import RiskLevel
from utils.color import (
    contrast_ratio,
    relative_luminance,
    required_ratio,
    suggest_accessible_color,
)

WHITE = (255, 255, 255)
BLACK = (0, 0, 0)


def _text_slide(fill: tuple[int, int, int] | None) -> tuple[Any, Any]:
    """Return a fresh presentation and a blank-layout slide with a textbox.

    When ``fill`` is given, the textbox gets a solid RGB background so the rule
    can resolve a deterministic background.
    """
    prs = Presentation()
    prs.core_properties.language = "en-US"
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(*fill)
    return prs, box


def _add_run(box: Any, text: str, *, color: tuple[int, int, int] | None = None, size_pt: float | None = None) -> Any:
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    if color is not None:
        run.font.color.rgb = RGBColor(*color)
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    return run


# --- pure color math -------------------------------------------------------


def test_contrast_ratio_extremes() -> None:
    """Black on white is the maximum 21:1; a color on itself is 1:1."""
    assert round(contrast_ratio(BLACK, WHITE), 1) == 21.0
    assert round(contrast_ratio(WHITE, WHITE), 1) == 1.0


def test_relative_luminance_bounds() -> None:
    """Luminance spans the full 0.0-1.0 range for black and white."""
    assert relative_luminance(BLACK) == 0.0
    assert round(relative_luminance(WHITE), 3) == 1.0


def test_suggest_accessible_color_reaches_target() -> None:
    """A failing gray-on-white pair is nudged to meet the AA normal ratio."""
    target = required_ratio(is_large=False)
    assert contrast_ratio((0xAA, 0xAA, 0xAA), WHITE) < target
    suggested = suggest_accessible_color((0xAA, 0xAA, 0xAA), WHITE, target=target)
    assert contrast_ratio(suggested, WHITE) >= target


def test_suggest_accessible_color_keeps_passing_color() -> None:
    """A color that already passes is returned unchanged."""
    assert suggest_accessible_color(BLACK, WHITE, target=4.5) == BLACK


# --- deterministic rule ----------------------------------------------------


def test_low_contrast_text_is_flagged() -> None:
    """Light gray text on a white fill is flagged as a medium-risk finding."""
    prs, box = _text_slide(fill=WHITE)
    _add_run(box, "Barely visible", color=(0xBB, 0xBB, 0xBB))

    findings = contrast.check(prs, Settings())

    assert len(findings) == 1
    finding = findings[0]
    assert finding.rule_id == "low_contrast_text"
    assert finding.wcag_criterion == "1.4.3 Contrast (Minimum)"
    assert finding.element_id == str(box.shape_id)
    assert finding.metadata["background"] == "#FFFFFF"
    assert finding.metadata["ratio"] < finding.metadata["required_ratio"]
    assert classify(finding) == RiskLevel.MEDIUM


def test_sufficient_contrast_is_not_flagged() -> None:
    """Black text on a white fill produces no finding."""
    prs, box = _text_slide(fill=WHITE)
    _add_run(box, "Clearly readable", color=BLACK)

    assert contrast.check(prs, Settings()) == []


def test_large_text_uses_lower_threshold() -> None:
    """A gray that fails at 4.5:1 passes the 3:1 large-text bar."""
    gray = (0x80, 0x80, 0x80)  # ~3.9:1 on white: fails normal, passes large
    assert contrast_ratio(gray, WHITE) < required_ratio(is_large=False)
    assert contrast_ratio(gray, WHITE) >= required_ratio(is_large=True)

    prs, box = _text_slide(fill=WHITE)
    _add_run(box, "Big heading", color=gray, size_pt=24)

    assert contrast.check(prs, Settings()) == []


def test_inherited_colors_on_default_theme_do_not_false_positive() -> None:
    """Unstyled text (no color, no fill) resolves to dark-on-light and passes.

    This is the key regression against the earlier false-positive flood: the
    default theme renders body text as dark on a light background, which the
    rule now measures instead of flagging.
    """
    prs, box = _text_slide(fill=None)
    _add_run(box, "Ordinary body text with no explicit color")

    assert contrast.check(prs, Settings()) == []


def test_theme_colored_text_is_resolved_and_measured() -> None:
    """A theme color is resolved via the palette, not guessed at.

    Setting the text to the theme's Background-1 (white) with an inherited
    (also white) background yields white-on-white — a real, deterministically
    detected violation, with no LLM involved.
    """
    prs, box = _text_slide(fill=None)
    run = _add_run(box, "Invisible theme text")
    run.font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1

    findings = contrast.check(prs, Settings())
    assert len(findings) == 1
    assert findings[0].rule_id == "low_contrast_text"
    assert findings[0].metadata["foreground"] == "#FFFFFF"
    assert findings[0].metadata["background"] == "#FFFFFF"


def test_theme_dark_text_on_light_theme_passes() -> None:
    """Theme Text-1 (dark) on the inherited light background produces no finding."""
    prs, box = _text_slide(fill=None)
    run = _add_run(box, "Readable theme text")
    run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1

    assert contrast.check(prs, Settings()) == []


def test_rule_runs_via_engine() -> None:
    """The rule is wired into the engine and its findings are triaged medium."""
    prs, box = _text_slide(fill=WHITE)
    _add_run(box, "Barely visible", color=(0xBB, 0xBB, 0xBB))

    findings = run_checks_on_presentation(prs, Settings())
    contrast_findings = [f for f in findings if f.rule_id == "low_contrast_text"]
    assert len(contrast_findings) == 1
    assert contrast_findings[0].risk_level == RiskLevel.MEDIUM
