"""Tests for the Tier-2 text-over-image/gradient contrast rule (LibreOffice-free).

Backgrounds are built with real solid-color PNGs (via Pillow) and real gradient
fills, so the rule samples genuine pixels and stops — no renderer, no network.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from checker.engine import run_checks_on_presentation
from checker.rules import contrast_image
from checker.triage import classify
from config import Settings
from models import RiskLevel


def _solid_png(path: Path, rgb: tuple[int, int, int]) -> None:
    Image.new("RGB", (48, 48), rgb).save(path)


def _deck_text_over_picture(
    tmp_path: Path, image_rgb: tuple[int, int, int], text_rgb: tuple[int, int, int] | None
) -> tuple[Any, Any]:
    """A blank slide with a full-bleed picture and a text box laid over it."""
    image_path = tmp_path / "bg.png"
    _solid_png(image_path, image_rgb)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(0), Inches(0), Inches(10), Inches(7.5))
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Overlay text"
    if text_rgb is not None:
        run.font.color.rgb = RGBColor(*text_rgb)
    return prs, box


def test_dark_text_over_dark_image_is_flagged(tmp_path: Path) -> None:
    """Low-contrast text sampled over a dark image is reported."""
    prs, box = _deck_text_over_picture(tmp_path, (20, 20, 20), (70, 70, 70))

    findings = contrast_image.check(prs, Settings())

    assert len(findings) == 1
    finding = findings[0]
    assert finding.rule_id == "low_contrast_over_image"
    assert finding.element_id == str(box.shape_id)
    assert finding.metadata["background_type"] == "image"
    assert finding.metadata["median_ratio"] < finding.metadata["required_ratio"]
    assert classify(finding) == RiskLevel.MEDIUM


def test_high_contrast_over_image_is_not_flagged(tmp_path: Path) -> None:
    """Black text over a white image passes and produces no finding."""
    prs, _ = _deck_text_over_picture(tmp_path, (255, 255, 255), (0, 0, 0))

    assert contrast_image.check(prs, Settings()) == []


def test_low_contrast_over_gradient_is_flagged(tmp_path: Path) -> None:
    """White text on a gradient whose stops are mid-tone fails somewhere and is flagged."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(6), Inches(1))
    shape.fill.gradient()
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "Banner text"
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    findings = contrast_image.check(prs, Settings())

    assert len(findings) == 1
    assert findings[0].rule_id == "low_contrast_over_gradient"
    assert findings[0].metadata["background_type"] == "gradient"
    assert classify(findings[0]) == RiskLevel.MEDIUM


def test_tier1_defers_so_no_double_report(tmp_path: Path) -> None:
    """Text over a picture is reported once (by the image rule), not also by Tier 1."""
    prs, box = _deck_text_over_picture(tmp_path, (20, 20, 20), (70, 70, 70))

    findings = run_checks_on_presentation(prs, Settings())
    rule_ids = [f.rule_id for f in findings if f.element_id == str(box.shape_id)]

    assert "low_contrast_over_image" in rule_ids
    assert "low_contrast_text" not in rule_ids


def test_no_visual_background_yields_nothing(tmp_path: Path) -> None:
    """A plain text box with no picture/gradient behind it is out of this rule's scope."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text_frame.paragraphs[0].add_run().text = "Just text"

    assert contrast_image.check(prs, Settings()) == []
