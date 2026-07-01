"""Tests for the optional, provider-pluggable LLM layer.

All tests inject a FakeProvider — no network calls and no API keys. They verify
that LLM output flows into suggestions/detection while the deterministic
pipeline and audit trail remain intact, and that the layer is fully gated off
by default.
"""

from __future__ import annotations

import base64
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from audit import ledger
from audit.rollback import rollback_item
from checker.engine import run_checks_on_presentation
from config import Settings
from fixer.auto_fix import apply_auto_fix
from fixer.flag_manual import record_manual_flag
from fixer.suggest_fix import approve_suggestion, record_pending_suggestion
from llm.factory import get_provider
from models import LedgerStatus, RiskLevel
from utils.pptx_xml import title_text

ONE_PIXEL_PNG = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p94AAAAASUVORK5CYII="

ALT_TEXT = "A bar chart of regional revenue growth."
GEN_TEXT = "Open the Q3 revenue report"


class FakeProvider:
    """In-memory LLMProvider used to drive the layer without a real SDK."""

    name = "fake:test-model"

    def __init__(self, *, text: str = GEN_TEXT, image: str = ALT_TEXT) -> None:
        self._text = text
        self._image = image

    def generate_text(self, *, system: str, prompt: str) -> str:
        return self._text

    def describe_image(self, *, image_bytes: bytes, media_type: str, system: str, prompt: str) -> str:
        return self._image


def _image_deck(tmp_path: Path) -> Path:
    image_path = tmp_path / "image001.png"
    image_path.write_bytes(base64.b64decode(ONE_PIXEL_PNG))
    pptx_path = tmp_path / "deck.pptx"
    prs = Presentation()
    prs.core_properties.language = "en-US"
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(0.5))
    title.text = "Regional revenue growth"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1.5), width=Inches(1))
    prs.save(pptx_path)
    return pptx_path


def _link_deck(tmp_path: Path) -> Path:
    pptx_path = tmp_path / "links.pptx"
    prs = Presentation()
    prs.core_properties.language = "en-US"
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(0.5))
    title.text = "Quarterly results"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "click here"
    run.hyperlink.address = "https://example.com/q3-report"
    prs.save(pptx_path)
    return pptx_path


def test_factory_returns_none_when_disabled_or_misconfigured() -> None:
    """The factory gates the layer off without erroring."""
    assert get_provider(Settings(llm_enabled=False)) is None
    assert get_provider(Settings(llm_enabled=True, llm_provider="openai", openai_api_key=None)) is None
    assert get_provider(Settings(llm_enabled=True, llm_provider="bogus", openai_api_key="k")) is None


def test_factory_builds_named_provider_without_network() -> None:
    """Selecting a provider with a key constructs it (no request is made)."""
    provider = get_provider(
        Settings(llm_enabled=True, llm_provider="openai", openai_api_key="sk-test", openai_model="gpt-4o-mini")
    )
    assert provider is not None
    assert provider.name == "openai:gpt-4o-mini"


def test_alt_text_uses_llm_when_provider_present(tmp_path: Path) -> None:
    """Image alt-text suggestions adopt the LLM output and record provenance."""
    pptx_path = _image_deck(tmp_path)
    ledger_path = tmp_path / "ledger.json"
    prs = Presentation(pptx_path)
    findings = run_checks_on_presentation(prs, Settings())
    alt_finding = next(f for f in findings if f.element_type == "image")
    entry = record_pending_suggestion(alt_finding, ledger_path, prs=prs, provider=FakeProvider())
    assert entry.suggested_fix == ALT_TEXT
    assert entry.metadata["suggestion_source"] == "fake:test-model"


def test_alt_text_falls_back_to_template_without_provider(tmp_path: Path) -> None:
    """Without a provider, the deterministic template suggestion is preserved."""
    pptx_path = _image_deck(tmp_path)
    ledger_path = tmp_path / "ledger.json"
    prs = Presentation(pptx_path)
    findings = run_checks_on_presentation(prs, Settings())
    alt_finding = next(f for f in findings if f.element_type == "image")
    entry = record_pending_suggestion(alt_finding, ledger_path)
    assert entry.suggested_fix is not None
    assert entry.suggested_fix != ALT_TEXT
    assert "suggestion_source" not in entry.metadata


def test_link_llm_autofix_is_flagged_and_reversible(tmp_path: Path) -> None:
    """LLM link text is auto-applied, flagged for review, and rolls back cleanly."""
    pptx_path = _link_deck(tmp_path)
    fixed_path = tmp_path / "fixed.pptx"
    rolled_path = tmp_path / "rolled.pptx"
    ledger_path = tmp_path / "ledger.json"

    prs = Presentation(pptx_path)
    findings = run_checks_on_presentation(prs, Settings())
    link_finding = next(f for f in findings if f.rule_id == "generic_link_text")
    entry = apply_auto_fix(prs, link_finding, ledger_path, Settings(), provider=FakeProvider())
    prs.save(fixed_path)

    assert entry is not None
    assert entry.status == LedgerStatus.AUTO_APPLIED
    assert entry.suggested_fix == GEN_TEXT
    assert entry.metadata["needs_review"] is True
    assert entry.snapshot_path

    rollback_item(pptx_path=fixed_path, output_path=rolled_path, ledger_path=ledger_path, item_id=entry.item_id)
    rolled = run_checks_on_presentation(Presentation(rolled_path), Settings())
    assert any(f.rule_id == "generic_link_text" for f in rolled)


def test_high_risk_draft_stored_without_editing_pptx(tmp_path: Path) -> None:
    """High-risk drafts land in metadata; suggested_fix and the PPTX are untouched."""
    pptx_path = tmp_path / "table.pptx"
    ledger_path = tmp_path / "ledger.json"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_table(10, 3, Inches(1), Inches(1), Inches(6), Inches(4))
    prs.save(pptx_path)

    prs = Presentation(pptx_path)
    findings = run_checks_on_presentation(prs, Settings())
    table_finding = next(f for f in findings if f.rule_id == "complex_table_review")
    assert table_finding.risk_level == RiskLevel.HIGH
    entry = record_manual_flag(table_finding, ledger_path, prs=prs, provider=FakeProvider())
    assert entry.status == LedgerStatus.FLAGGED_MANUAL
    assert entry.suggested_fix is None
    assert entry.metadata["llm_draft"] == GEN_TEXT


def test_detection_pass_appends_findings_with_provenance(tmp_path: Path) -> None:
    """The semantic detection pass adds findings tagged with their provider."""
    pptx_path = _link_deck(tmp_path)
    prs = Presentation(pptx_path)
    settings = Settings(llm_enabled=True)
    findings = run_checks_on_presentation(prs, settings, provider=FakeProvider(text="Title is too vague"))
    llm_findings = [f for f in findings if f.rule_id == "llm_semantic_review"]
    assert llm_findings
    assert all(f.risk_level == RiskLevel.MEDIUM for f in llm_findings)
    assert all(f.metadata["detected_by"] == "fake:test-model" for f in llm_findings)


def test_weak_title_suggestion_applies_and_rolls_back(tmp_path: Path) -> None:
    """A vague title is detected, suggested, applied on approval, and reversible."""
    pptx_path = tmp_path / "weak.pptx"
    prs = Presentation()
    prs.core_properties.language = "en-US"
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide 2"
    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    body.text_frame.text = "Quarterly revenue grew across all regions."
    prs.save(pptx_path)

    ledger_path = tmp_path / "ledger.json"
    fixed_path = tmp_path / "fixed.pptx"
    rolled_path = tmp_path / "rolled.pptx"
    better = "Regional Revenue Growth"

    prs = Presentation(pptx_path)
    findings = run_checks_on_presentation(prs, Settings(llm_enabled=True), provider=FakeProvider(text=better))
    title_finding = next(f for f in findings if f.rule_id == "weak_slide_title")
    assert title_finding.risk_level == RiskLevel.MEDIUM
    assert title_finding.suggested_fix == better

    entry = record_pending_suggestion(title_finding, ledger_path, prs=prs, provider=FakeProvider(text=better))
    assert entry.status == LedgerStatus.PENDING_APPROVAL

    approved = approve_suggestion(prs=prs, ledger_path=ledger_path, item_id=entry.item_id, approved_by="alice")
    assert approved.status == LedgerStatus.APPROVED
    prs.save(fixed_path)
    assert title_text(Presentation(fixed_path).slides[0]) == better

    rollback_item(pptx_path=fixed_path, output_path=rolled_path, ledger_path=ledger_path, item_id=entry.item_id)
    assert title_text(Presentation(rolled_path).slides[0]) == "Slide 2"


def test_detection_skipped_when_disabled(tmp_path: Path) -> None:
    """No LLM findings appear when the layer is disabled, even if a provider is passed."""
    pptx_path = _link_deck(tmp_path)
    prs = Presentation(pptx_path)
    findings = run_checks_on_presentation(prs, Settings(llm_enabled=False), provider=FakeProvider())
    assert not any(f.rule_id == "llm_semantic_review" for f in findings)
