"""Smoke tests for the essential accessibility-agent flow."""

from __future__ import annotations

import base64
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from audit import ledger
from audit.rollback import rollback_item
from checker.engine import run_checks_on_presentation
from config import Settings
from fixer.auto_fix import apply_auto_fix
from fixer.suggest_fix import approve_suggestion, record_pending_suggestion
from models import RiskLevel
from pptx.oxml.ns import qn
from utils.pptx_xml import get_alt_text, get_shape_by_id

ONE_PIXEL_PNG = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p94AAAAASUVORK5CYII="


def _group_shapes(slide: object, shapes: list) -> None:
    """Move existing top-level shapes into a single PowerPoint group."""
    sp_tree = slide.shapes._spTree
    group = sp_tree.makeelement(qn("p:grpSp"), {})
    nv = group.makeelement(qn("p:nvGrpSpPr"), {})
    nv.append(group.makeelement(qn("p:cNvPr"), {"id": "900", "name": "FigureGroup"}))
    nv.append(group.makeelement(qn("p:cNvGrpSpPr"), {}))
    nv.append(group.makeelement(qn("p:nvPr"), {}))
    group.append(nv)
    grp_pr = group.makeelement(qn("p:grpSpPr"), {})
    xfrm = grp_pr.makeelement(qn("a:xfrm"), {})
    for tag, attrs in (
        ("a:off", {"x": "0", "y": "0"}),
        ("a:ext", {"cx": "1000000", "cy": "1000000"}),
        ("a:chOff", {"x": "0", "y": "0"}),
        ("a:chExt", {"cx": "1000000", "cy": "1000000"}),
    ):
        xfrm.append(xfrm.makeelement(qn(tag), attrs))
    grp_pr.append(xfrm)
    group.append(grp_pr)
    for shape in shapes:
        sp_tree.remove(shape._element)
        group.append(shape._element)
    sp_tree.append(group)


def _save_link_deck(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "click here"
    run.hyperlink.address = "https://example.com/resource"
    prs.save(path)


def test_low_risk_auto_fix_and_rollback(tmp_path: Path) -> None:
    """Low-risk findings are fixed, logged, and individually reversible."""
    input_path = tmp_path / "input.pptx"
    fixed_path = tmp_path / "fixed.pptx"
    rollback_path = tmp_path / "rollback.pptx"
    ledger_path = tmp_path / "ledger.json"
    _save_link_deck(input_path)

    settings = Settings()
    prs = Presentation(input_path)
    findings = run_checks_on_presentation(prs, settings)
    for finding in findings:
        if finding.risk_level == RiskLevel.LOW:
            apply_auto_fix(prs, finding, ledger_path, settings)
    prs.save(fixed_path)

    entries = ledger.load_ledger(ledger_path)
    assert len(entries) == 3
    generic_link_entry = next(entry for entry in entries if entry.rule_id == "generic_link_text")
    rollback_item(
        pptx_path=fixed_path,
        output_path=rollback_path,
        ledger_path=ledger_path,
        item_id=generic_link_entry.item_id,
    )
    rolled_prs = Presentation(rollback_path)
    rolled_findings = run_checks_on_presentation(rolled_prs, settings)
    assert any(finding.rule_id == "generic_link_text" for finding in rolled_findings)


def test_medium_alt_text_suggestion_can_be_approved(tmp_path: Path) -> None:
    """Medium-risk image alt text suggestions are staged and approved later."""
    image_path = tmp_path / "image001.png"
    image_path.write_bytes(base64.b64decode(ONE_PIXEL_PNG))
    pptx_path = tmp_path / "image_deck.pptx"
    ledger_path = tmp_path / "ledger.json"
    reviewed_path = tmp_path / "reviewed.pptx"

    prs = Presentation()
    prs.core_properties.language = "en-US"
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(0.5))
    title.text = "Regional revenue growth"
    picture = slide.shapes.add_picture(str(image_path), Inches(1), Inches(1.5), width=Inches(1))
    prs.save(pptx_path)

    prs = Presentation(pptx_path)
    findings = run_checks_on_presentation(prs, Settings())
    alt_finding = next(finding for finding in findings if finding.rule_id == "weak_image_alt_text")
    entry = record_pending_suggestion(alt_finding, ledger_path)
    approve_suggestion(prs=prs, ledger_path=ledger_path, item_id=entry.item_id, approved_by="test")
    prs.save(reviewed_path)

    reviewed = Presentation(reviewed_path)
    reviewed_picture = next(shape for shape in reviewed.slides[0].shapes if shape.shape_id == picture.shape_id)
    assert get_alt_text(reviewed_picture).startswith("Image related to")


def test_multiple_grouped_figures_each_get_alt_text(tmp_path: Path) -> None:
    """Every picture inside a group is detected and fixable, not just the first."""
    image_path = tmp_path / "image001.png"
    image_path.write_bytes(base64.b64decode(ONE_PIXEL_PNG))
    pptx_path = tmp_path / "grouped_deck.pptx"
    ledger_path = tmp_path / "ledger.json"
    reviewed_path = tmp_path / "reviewed.pptx"

    prs = Presentation()
    prs.core_properties.language = "en-US"
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.5))
    title.text = "Quarterly charts"
    pic_a = slide.shapes.add_picture(str(image_path), Inches(1), Inches(1.5), width=Inches(1))
    pic_b = slide.shapes.add_picture(str(image_path), Inches(3), Inches(1.5), width=Inches(1))
    _group_shapes(slide, [pic_a, pic_b])
    prs.save(pptx_path)

    prs = Presentation(pptx_path)
    findings = run_checks_on_presentation(prs, Settings())
    alt_findings = [f for f in findings if f.element_type == "image"]
    assert {f.element_id for f in alt_findings} == {str(pic_a.shape_id), str(pic_b.shape_id)}

    for finding in alt_findings:
        entry = record_pending_suggestion(finding, ledger_path)
        approve_suggestion(prs=prs, ledger_path=ledger_path, item_id=entry.item_id, approved_by="test")
    prs.save(reviewed_path)

    reviewed = Presentation(reviewed_path)
    for pic in (pic_a, pic_b):
        shape = get_shape_by_id(reviewed.slides[0], pic.shape_id)
        assert get_alt_text(shape).startswith("Image related to")
