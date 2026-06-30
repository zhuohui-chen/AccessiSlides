"""Low-risk automatic remediation engine."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx.util import Inches, Pt

from audit import ledger
from audit.snapshot import (
    create_added_shape_snapshot,
    create_core_property_snapshot,
    create_shape_snapshot,
)
from config import Settings
from models import Finding, LedgerEntry, LedgerStatus, RiskLevel
from utils.logging import get_logger
from utils.pptx_xml import get_shape_by_id, get_slide, slide_context_text

LOGGER = get_logger(__name__)


def _apply_llm_text(finding: Finding, suggestion: str | None, provider: Any) -> None:
    """Adopt an LLM suggestion as the finding's fix and flag it for later review.

    The deterministic text remains the fallback when ``suggestion`` is empty.
    Items enhanced here stay AUTO_APPLIED but carry ``needs_review`` so the
    interactive report surfaces them for human confirmation.
    """
    if not suggestion:
        return
    finding.suggested_fix = suggestion
    finding.metadata = {
        **finding.metadata,
        "suggestion_source": provider.name,
        "needs_review": True,
    }


def _add_slide_title(
    prs: Any, finding: Finding, ledger_path: Path, settings: Settings, provider: Any | None = None
) -> LedgerEntry:
    """Add a visible title to a slide that has none.

    With a provider, the title text is generated from the slide content and the
    item is flagged for later human review; otherwise a deterministic
    ``"<prefix> <n>"`` placeholder is used.
    """
    item_id = ledger.next_item_id(ledger_path)
    slide = get_slide(prs, finding.slide_number)
    if provider is not None:
        from llm import service

        _apply_llm_text(
            finding,
            service.suggest_slide_title(provider, context=slide_context_text(slide)),
            provider,
        )
    title_text = finding.suggested_fix or f"{settings.auto_title_prefix} {finding.slide_number}"
    left = Inches(0.25)
    top = Inches(0.12)
    width = max(prs.slide_width - Inches(0.5), Inches(1.0))
    height = Inches(0.35)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.name = "Auto Accessibility Title"
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.text = title_text
    if paragraph.runs:
        paragraph.runs[0].font.size = Pt(16)

    snapshot_path = create_added_shape_snapshot(
        ledger_path=ledger_path,
        item_id=item_id,
        slide_number=finding.slide_number,
        created_shape_id=str(textbox.shape_id),
        settings=settings,
    )
    entry = LedgerEntry.from_finding(
        item_id=item_id,
        finding=finding,
        status=LedgerStatus.AUTO_APPLIED,
        snapshot_path=snapshot_path,
        applied_at=ledger.utc_now(),
        approved_by="auto",
    )
    entry.element_id = str(textbox.shape_id)
    entry.metadata = {**entry.metadata, "created_shape_id": str(textbox.shape_id)}
    ledger.append_entry(ledger_path, entry)
    LOGGER.info("auto_fix_applied", item_id=item_id, rule_id=finding.rule_id)
    return entry


def _set_presentation_language(
    prs: Any, finding: Finding, ledger_path: Path, settings: Settings, provider: Any | None = None
) -> LedgerEntry:
    """Set missing presentation language metadata (deterministic, never LLM)."""
    del provider
    item_id = ledger.next_item_id(ledger_path)
    snapshot_path = create_core_property_snapshot(
        prs=prs,
        ledger_path=ledger_path,
        item_id=item_id,
        property_name="language",
        settings=settings,
    )
    prs.core_properties.language = finding.suggested_fix or settings.default_language
    entry = LedgerEntry.from_finding(
        item_id=item_id,
        finding=finding,
        status=LedgerStatus.AUTO_APPLIED,
        snapshot_path=snapshot_path,
        applied_at=ledger.utc_now(),
        approved_by="auto",
    )
    ledger.append_entry(ledger_path, entry)
    LOGGER.info("auto_fix_applied", item_id=item_id, rule_id=finding.rule_id)
    return entry


def _replace_generic_link_text(
    prs: Any, finding: Finding, ledger_path: Path, settings: Settings, provider: Any | None = None
) -> LedgerEntry:
    """Replace ambiguous hyperlink text with descriptive text.

    With a provider, the replacement is generated from the link target and slide
    context and the item is flagged for later human review; otherwise the
    deterministic URL-derived text is used.
    """
    item_id = ledger.next_item_id(ledger_path)
    snapshot_path = create_shape_snapshot(
        prs=prs,
        ledger_path=ledger_path,
        item_id=item_id,
        slide_number=finding.slide_number,
        shape_id=finding.element_id,
        settings=settings,
    )
    slide = get_slide(prs, finding.slide_number)
    shape = get_shape_by_id(slide, finding.element_id)
    paragraph_index = int(finding.metadata.get("paragraph_index", 0))
    run_index = int(finding.metadata.get("run_index", 0))
    if provider is not None:
        from llm import service

        _apply_llm_text(
            finding,
            service.suggest_link_text(
                provider,
                address=str(finding.metadata.get("address", "")),
                context=slide_context_text(slide, exclude_shape_id=finding.element_id),
            ),
            provider,
        )
    replacement = finding.suggested_fix or "Open linked resource"
    try:
        run = shape.text_frame.paragraphs[paragraph_index].runs[run_index]
    except (AttributeError, IndexError):
        raise ValueError(f"Cannot locate hyperlink run for finding {finding.rule_id}") from None
    run.text = replacement

    entry = LedgerEntry.from_finding(
        item_id=item_id,
        finding=finding,
        status=LedgerStatus.AUTO_APPLIED,
        snapshot_path=snapshot_path,
        applied_at=ledger.utc_now(),
        approved_by="auto",
    )
    ledger.append_entry(ledger_path, entry)
    LOGGER.info("auto_fix_applied", item_id=item_id, rule_id=finding.rule_id)
    return entry


def apply_auto_fix(
    prs: Any,
    finding: Finding,
    ledger_path: Path,
    settings: Settings | None = None,
    *,
    provider: Any | None = None,
) -> LedgerEntry | None:
    """Apply one low-risk fix and write its ledger entry.

    Args:
        prs: Loaded presentation to mutate.
        finding: Low-risk checker finding.
        ledger_path: Ledger output path.
        settings: Runtime settings.
        provider: Optional LLM provider. When supplied, slide-title and link-text
            fixes use generated text (flagged ``needs_review``); language metadata
            stays deterministic regardless.

    Returns:
        The ledger entry, or None when the rule has no auto-fixer and strict
        unknown-fix handling is disabled.
    """
    resolved_settings = settings or Settings()
    if finding.risk_level != RiskLevel.LOW:
        raise ValueError(f"Auto-fix only accepts low-risk findings, got {finding.risk_level.value}")
    handlers = {
        "missing_slide_title": _add_slide_title,
        "missing_presentation_language": _set_presentation_language,
        "generic_link_text": _replace_generic_link_text,
    }
    handler = handlers.get(finding.rule_id)
    if handler is None:
        if resolved_settings.fail_on_unknown_auto_fix:
            raise ValueError(f"No auto-fix handler for rule {finding.rule_id}")
        return None
    return handler(prs, finding, ledger_path, resolved_settings, provider)
