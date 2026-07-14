"""Low-level python-pptx XML helpers.

The public python-pptx API does not expose every accessibility-related field.
PowerPoint stores image alt text in the `p:cNvPr` XML element as the `descr`
attribute, so a few targeted XML helpers are needed.
"""

from __future__ import annotations

from collections.abc import Iterable
from typing import Any

from lxml import etree
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml


def _iter_elements_with_local_name(root: Any, local_name: str) -> Iterable[Any]:
    """Yield descendants whose tag has the requested XML local name."""
    for element in root.iter():
        if etree.QName(element).localname == local_name:
            yield element


def c_nv_pr(shape: Any) -> Any | None:
    """Return the non-visual properties element for a shape, if present."""
    return next(_iter_elements_with_local_name(shape._element, "cNvPr"), None)


def get_shape_name(shape: Any) -> str:
    """Return the PowerPoint shape name."""
    element = c_nv_pr(shape)
    return "" if element is None else str(element.get("name", ""))


def get_alt_text(shape: Any) -> str:
    """Return a shape's alt text description from the underlying XML."""
    element = c_nv_pr(shape)
    if element is None:
        return ""
    return str(element.get("descr") or "").strip()


def set_alt_text(shape: Any, alt_text: str) -> None:
    """Set a shape's alt text description in the underlying XML."""
    element = c_nv_pr(shape)
    if element is None:
        raise ValueError("Shape does not contain a p:cNvPr element")
    element.set("descr", alt_text.strip())


def shape_xml(shape: Any) -> str:
    """Serialize a shape XML element."""
    return etree.tostring(shape._element, encoding="unicode")


def replace_shape_xml(shape: Any, xml: str) -> None:
    """Replace an existing shape XML element with a serialized snapshot."""
    parent = shape._element.getparent()
    if parent is None:
        raise ValueError("Cannot replace shape XML because it has no parent")
    parent.replace(shape._element, parse_xml(xml))


def get_image_blob(shape: Any) -> tuple[bytes, str] | None:
    """Return a picture shape's raw bytes and MIME type, or None.

    Args:
        shape: A python-pptx shape expected to be a picture.

    Returns:
        A ``(blob, content_type)`` tuple, or ``None`` when the shape has no
        accessible image (e.g. it is not a picture, or its blob cannot be read).
    """
    image = getattr(shape, "image", None)
    if image is None:
        return None
    try:
        blob = image.blob
        content_type = image.content_type
    except (AttributeError, ValueError, KeyError):
        return None
    if not blob or not content_type:
        return None
    return blob, str(content_type)


# --- Theme color resolution -------------------------------------------------
#
# A run/shape color may be a literal RGB, a *theme* (scheme) color, or unset
# (inherited). WCAG 1.4.3 is about the *rendered* color, so we resolve theme and
# inherited colors to concrete RGB using the slide master's theme palette rather
# than guessing. This lets the contrast rule measure ordinary decks (whose text
# inherits the theme's dark text on a light background) instead of flagging them.

# MSO_THEME_COLOR member name -> theme ``clrScheme`` slot (logical names tx*/bg*
# are further mapped through the master's ``clrMap`` in :func:`slide_theme_palette`).
_THEME_SLOTS = {
    "DARK_1": "dk1", "LIGHT_1": "lt1", "DARK_2": "dk2", "LIGHT_2": "lt2",
    "ACCENT_1": "accent1", "ACCENT_2": "accent2", "ACCENT_3": "accent3",
    "ACCENT_4": "accent4", "ACCENT_5": "accent5", "ACCENT_6": "accent6",
    "HYPERLINK": "hlink", "FOLLOWED_HYPERLINK": "folHlink",
    "TEXT_1": "tx1", "TEXT_2": "tx2", "BACKGROUND_1": "bg1", "BACKGROUND_2": "bg2",
}
# Fallback RGB for ``sysClr`` entries that omit a ``lastClr`` attribute.
_SYS_COLOR_DEFAULTS = {"windowText": "000000", "window": "FFFFFF"}


def _hex_to_rgb(value: str) -> tuple[int, int, int] | None:
    """Parse an ``RRGGBB`` hex string into an RGB triple, or None."""
    text = (value or "").strip().lstrip("#")
    if len(text) != 6:
        return None
    try:
        return (int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))
    except ValueError:
        return None


def _scheme_slot_rgb(slot_element: Any) -> tuple[int, int, int] | None:
    """Return the RGB of a ``clrScheme`` slot (``srgbClr`` or ``sysClr``)."""
    for child in slot_element:
        local = etree.QName(child).localname
        if local == "srgbClr":
            return _hex_to_rgb(child.get("val", ""))
        if local == "sysClr":
            last = child.get("lastClr")
            if last:
                return _hex_to_rgb(last)
            return _hex_to_rgb(_SYS_COLOR_DEFAULTS.get(child.get("val", ""), ""))
    return None


def slide_theme_palette(slide: Any) -> dict[str, tuple[int, int, int]]:
    """Return the slide master's theme palette as ``slot -> RGB``.

    Keys include the raw scheme slots (``dk1``, ``lt1`` … ``accent6``) plus the
    logical names ``tx1``/``tx2``/``bg1``/``bg2`` resolved through the master's
    color map. Returns an empty dict when the theme cannot be read, in which case
    the contrast rule simply skips runs it cannot measure.
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    master = getattr(getattr(slide, "slide_layout", None), "slide_master", None)
    if master is None:
        return {}
    try:
        theme_part = master.part.part_related_by(RT.THEME)
        root = etree.fromstring(theme_part.blob)
    except (KeyError, AttributeError, etree.XMLSyntaxError):
        return {}
    scheme = next(_iter_elements_with_local_name(root, "clrScheme"), None)
    if scheme is None:
        return {}
    palette: dict[str, tuple[int, int, int]] = {}
    for slot in scheme:
        rgb = _scheme_slot_rgb(slot)
        if rgb is not None:
            palette[etree.QName(slot).localname] = rgb
    master_element = getattr(master, "element", None)
    color_map = (
        next(_iter_elements_with_local_name(master_element, "clrMap"), None)
        if master_element is not None
        else None
    )
    default_map = {"tx1": "dk1", "bg1": "lt1", "tx2": "dk2", "bg2": "lt2"}
    for logical, fallback in default_map.items():
        mapped = color_map.get(logical) if color_map is not None else fallback
        if mapped and mapped in palette:
            palette[logical] = palette[mapped]
    return palette


def _apply_brightness(rgb: tuple[int, int, int], brightness: float) -> tuple[int, int, int]:
    """Lighten (toward white) or darken (toward black) an RGB triple.

    Approximates PowerPoint's "Lighter/Darker" theme-color variants so that,
    e.g., "Text 1, Lighter 50%" resolves to a genuinely lighter gray.
    """
    if not brightness:
        return rgb
    if brightness > 0:
        return tuple(round(c + (255 - c) * brightness) for c in rgb)  # type: ignore[return-value]
    return tuple(round(c * (1.0 + brightness)) for c in rgb)  # type: ignore[return-value]


def resolve_color_rgb(
    color: Any, palette: dict[str, tuple[int, int, int]]
) -> tuple[int, int, int] | None:
    """Resolve a ``ColorFormat`` to concrete RGB using the theme ``palette``.

    Handles literal RGB and theme/scheme colors (including their brightness
    variant). Returns ``None`` for inherited/unset colors or theme colors the
    palette does not define, letting the caller apply a sensible default.
    """
    if color is None:
        return None
    try:
        color_type = color.type
    except (AttributeError, TypeError):
        return None
    if color_type is None:
        return None
    base: tuple[int, int, int] | None = None
    if color_type == MSO_COLOR_TYPE.RGB:
        try:
            rgb = color.rgb
            base = (rgb[0], rgb[1], rgb[2])
        except (AttributeError, TypeError):
            return None
    elif color_type == MSO_COLOR_TYPE.SCHEME:
        try:
            theme_color = color.theme_color
        except (AttributeError, TypeError):
            return None
        base = palette.get(_THEME_SLOTS.get(getattr(theme_color, "name", ""), ""))
        if base is None:
            return None
    else:
        return None
    try:
        brightness = color.brightness or 0.0
    except (AttributeError, TypeError):
        brightness = 0.0
    result = _apply_brightness(base, brightness)
    return tuple(max(0, min(255, int(c))) for c in result)  # type: ignore[return-value]


def resolve_run_rgb(
    run: Any, palette: dict[str, tuple[int, int, int]]
) -> tuple[int, int, int] | None:
    """Resolve a text run's font color to RGB, or None when inherited/unset."""
    font = getattr(run, "font", None)
    if font is None:
        return None
    return resolve_color_rgb(font.color, palette)


def _fill_rgb(fill: Any, palette: dict[str, tuple[int, int, int]]) -> tuple[int, int, int] | None:
    """Return the resolved RGB of a solid ``FillFormat``, or None otherwise."""
    if fill is None:
        return None
    try:
        if fill.type != MSO_FILL_TYPE.SOLID:
            return None
        fore_color = fill.fore_color
    except (AttributeError, TypeError, ValueError):
        return None
    return resolve_color_rgb(fore_color, palette)


def shape_background_rgb(
    shape: Any, slide: Any, palette: dict[str, tuple[int, int, int]]
) -> tuple[int, int, int] | None:
    """Resolve the effective background RGB behind a shape's text, or None.

    Checks, in order: the shape's own solid fill, the slide background, the
    layout background, and the master background — each resolved through the
    theme ``palette``. Returns ``None`` when none resolve; the caller then falls
    back to the theme's default background color.
    """
    own = _fill_rgb(getattr(shape, "fill", None), palette)
    if own is not None:
        return own
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None)
    for source in (slide, layout, master):
        if source is None:
            continue
        background = getattr(source, "background", None)
        rgb = _fill_rgb(getattr(background, "fill", None), palette) if background is not None else None
        if rgb is not None:
            return rgb
    return None


# --- Geometry & visual-background detection (for text-over-image contrast) ---

_LARGE_PT = 18.0
_LARGE_BOLD_PT = 14.0


def run_is_large_text(run: Any) -> bool:
    """Return True when a run qualifies as WCAG large text (lower contrast bar)."""
    font = getattr(run, "font", None)
    if font is None:
        return False
    size = font.size
    if size is None:
        return False
    point_size = size.pt
    if bool(font.bold) and point_size >= _LARGE_BOLD_PT:
        return True
    return point_size >= _LARGE_PT


def shape_rect(shape: Any) -> tuple[int, int, int, int] | None:
    """Return a shape's ``(left, top, width, height)`` in EMU, or None if unset."""
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    except AttributeError:
        return None
    if left is None or top is None or width is None or height is None:
        return None
    if int(width) <= 0 or int(height) <= 0:
        return None
    return (int(left), int(top), int(width), int(height))


def rect_intersection(
    a: tuple[int, int, int, int], b: tuple[int, int, int, int]
) -> tuple[int, int, int, int] | None:
    """Return the overlap ``(left, top, width, height)`` of two rects, or None."""
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    x0, y0 = max(ax, bx), max(ay, by)
    x1, y1 = min(ax + aw, bx + bw), min(ay + ah, by + bh)
    if x1 <= x0 or y1 <= y0:
        return None
    return (x0, y0, x1 - x0, y1 - y0)


def overlap_fraction(
    text_rect: tuple[int, int, int, int], other_rect: tuple[int, int, int, int]
) -> float:
    """Return the fraction of ``text_rect``'s area covered by ``other_rect``."""
    inter = rect_intersection(text_rect, other_rect)
    if inter is None:
        return 0.0
    text_area = text_rect[2] * text_rect[3]
    if text_area <= 0:
        return 0.0
    return (inter[2] * inter[3]) / text_area


def is_picture(shape: Any) -> bool:
    """Return True when a shape is a picture."""
    return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE


def shape_has_gradient_fill(shape: Any) -> bool:
    """Return True when a shape has a gradient fill."""
    try:
        return shape.fill.type == MSO_FILL_TYPE.GRADIENT
    except (AttributeError, TypeError, ValueError):
        return False


def gradient_stops_rgb(
    shape: Any, palette: dict[str, tuple[int, int, int]]
) -> list[tuple[int, int, int]]:
    """Return the resolved RGB of a shape's gradient stops, or [] when not applicable."""
    try:
        fill = shape.fill
        if fill.type != MSO_FILL_TYPE.GRADIENT:
            return []
        stops = list(fill.gradient_stops)
    except (AttributeError, TypeError, ValueError, NotImplementedError):
        return []
    resolved: list[tuple[int, int, int]] = []
    for stop in stops:
        try:
            rgb = resolve_color_rgb(stop.color, palette)
        except (AttributeError, TypeError, ValueError):
            continue
        if rgb is not None:
            resolved.append(rgb)
    return resolved


def find_visual_background(
    text_shape: Any, slide: Any, *, min_overlap: float = 0.5
) -> tuple[str, Any] | None:
    """Return the picture/gradient background beneath a text shape, or None.

    Detects when a text shape either carries its own gradient fill or overlaps an
    underlying (lower z-order) picture or gradient-filled shape by at least
    ``min_overlap`` of the text shape's area. Returns ``(kind, source_shape)``
    where ``kind`` is ``"picture"`` or ``"gradient"``.

    Only top-level slide shapes participate in the overlap search (grouped shapes
    use relative coordinates); a grouped text shape can still match on its own
    gradient fill. This is the seam the deterministic XML rule uses to defer such
    runs to the pixel/gradient-sampling rule, so the two never double-report.
    """
    if shape_has_gradient_fill(text_shape):
        return ("gradient", text_shape)
    text_rect = shape_rect(text_shape)
    if text_rect is None:
        return None
    found: tuple[str, Any] | None = None
    seen_text = False
    for other in slide.shapes:
        # python-pptx yields a fresh proxy per iteration, so compare the shared
        # underlying element rather than object identity.
        if other._element is text_shape._element:
            seen_text = True
            break
        other_rect = shape_rect(other)
        if other_rect is None or overlap_fraction(text_rect, other_rect) < min_overlap:
            continue
        if is_picture(other):
            found = ("picture", other)
        elif shape_has_gradient_fill(other):
            found = ("gradient", other)
    return found if seen_text else None


def get_slide(prs: Any, slide_number: int) -> Any:
    """Return a 1-indexed slide from a presentation."""
    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide number out of range: {slide_number}")
    return prs.slides[slide_number - 1]


def iter_shapes(shapes: Any) -> Iterable[Any]:
    """Yield every shape in a container, descending into group shapes.

    PowerPoint groups (``MSO_SHAPE_TYPE.GROUP``) nest their members, and
    iterating a slide's shape tree does not recurse into them. This walker
    yields each group shape and then its descendants, so a slide that groups
    several charts or pictures exposes every figure individually.
    """
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def get_shape_by_id(slide: Any, shape_id: str | int) -> Any:
    """Return a shape by PowerPoint shape ID, searching inside groups."""
    target = str(shape_id)
    for shape in iter_shapes(slide.shapes):
        if str(shape.shape_id) == target:
            return shape
    raise KeyError(f"Shape id {target} not found on slide")


def delete_shape_by_id(slide: Any, shape_id: str | int) -> None:
    """Remove a shape by PowerPoint shape ID."""
    shape = get_shape_by_id(slide, shape_id)
    parent = shape._element.getparent()
    if parent is None:
        raise ValueError("Cannot delete shape XML because it has no parent")
    parent.remove(shape._element)


def text_from_shape(shape: Any) -> str:
    """Return text for a shape when it has a text frame."""
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                parts.append(run.text)
        if paragraph.text and not paragraph.runs:
            parts.append(paragraph.text)
    return " ".join(part.strip() for part in parts if part.strip())


def slide_context_text(slide: Any, *, exclude_shape_id: str | None = None, limit: int = 180) -> str:
    """Return concise text context from a slide."""
    snippets: list[str] = []
    for shape in iter_shapes(slide.shapes):
        if exclude_shape_id is not None and str(shape.shape_id) == str(exclude_shape_id):
            continue
        text = text_from_shape(shape).strip()
        if text:
            snippets.append(text)
    context = " ".join(snippets)
    return context[:limit].strip()


def title_shape(slide: Any) -> Any | None:
    """Return the shape holding the slide's title text, if one carries text.

    Prefers the layout title placeholder, then falls back to any shape whose
    name contains "title". Returns ``None`` when no title-bearing shape exists.
    """
    placeholder = getattr(slide.shapes, "title", None)
    if placeholder is not None and text_from_shape(placeholder).strip():
        return placeholder
    for shape in slide.shapes:
        if "title" in get_shape_name(shape).lower() and text_from_shape(shape).strip():
            return shape
    return None


def title_text(slide: Any) -> str:
    """Return the best available slide title text."""
    shape = title_shape(slide)
    return "" if shape is None else text_from_shape(shape).strip()


def set_title_text(shape: Any, text: str) -> None:
    """Replace a title shape's text with ``text``, keeping its first run's font."""
    if not getattr(shape, "has_text_frame", False):
        raise ValueError("Title shape has no text frame")
    text_frame = shape.text_frame
    paragraph = text_frame.paragraphs[0]
    font = paragraph.runs[0].font if paragraph.runs else None
    size = font.size if font is not None else None
    text_frame.text = text.strip()
    new_paragraph = text_frame.paragraphs[0]
    if size is not None and new_paragraph.runs:
        new_paragraph.runs[0].font.size = size
